# -*- coding: utf-8 -*-

"""Classification_Test

Automatically generated script for local execution.

### **Install modules**
"""

# Install required modules
import subprocess
import sys

# required_modules = [
#     "txtai",
#     "datasets",
#     "transformers",
#     "python-docx",
#     "python-pptx",
#     "pandas",
#     "docx2txt",
#     "PyPDF2",
#     "textract"
# ]

# for module in required_modules:
#     subprocess.check_call([sys.executable, "-m", "pip", "install", module])

# """Function to read **files**

# """

from flask import Flask, request, render_template, redirect, url_for
import os
import io
from txtai.embeddings import Embeddings
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploaded_files'

# Initialize embeddings model
embeddings = Embeddings({"path": "sentence-transformers/paraphrase-MiniLM-L3-v2", "content": True})

# Ensure upload directory exists
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

def process_files(upload_dir):
    data = []
    for file in os.listdir(upload_dir):
        file_path = os.path.join(upload_dir, file)
        with open(file_path, 'rb') as f:
            content = f.read()
        data.extend(process_file_content(file, content))
    return data

def process_file_content(file_name, content):
    if file_name.endswith(".txt"):
        return process_text_content(content)
    elif file_name.endswith(".csv"):
        return process_csv_content(content)
    elif file_name.endswith(".pdf"):
        return process_pdf_content(content)
    elif file_name.endswith(".docx"):
        return process_docx_content(content)
    elif file_name.endswith(".pptx"):
        return process_pptx_content(content)
    return []

def process_text_content(content):
    text = content.decode("utf-8")
    paragraphs = text.split("\n\n")
    return [{"text": para.strip()} for para in paragraphs if para.strip()]

def process_csv_content(content):
    text = content.decode("utf-8")
    lines = text.splitlines()
    return [{"text": line.strip()} for line in lines if line.strip()]

def process_pdf_content(content):
    custom_data = []
    pdf_reader = PdfReader(io.BytesIO(content))
    for page_num in range(len(pdf_reader.pages)):
        page = pdf_reader.pages[page_num]
        text = page.extract_text()
        if text:
            paragraphs = text.split("\n\n")
            custom_data.extend([{"text": para.strip()} for para in paragraphs if para.strip()])
    return custom_data

def process_docx_content(content):
    doc = Document(io.BytesIO(content))
    custom_data = []
    for para in doc.paragraphs:
        if para.text.strip():
            custom_data.append({"text": para.text.strip()})
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    custom_data.append({"text": cell.text.strip()})
    return custom_data

def process_pptx_content(content):
    custom_data = []
    presentation = Presentation(io.BytesIO(content))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    paragraphs = text.split("\n\n")
                    custom_data.extend([{"text": para.strip()} for para in paragraphs if para.strip()])
    return custom_data

def create_embeddings(data):
    embeddings.index(data)

def search(query):
    results = embeddings.search(query, limit=100 )  # Adjust limit as needed
    return results

def format_search_results(results):
    formatted_results = []
    for idx, result in enumerate(results, start=1):
        formatted_results.append(f"Result {idx}:")
        formatted_results.append(f"Score: {result['score']:.4f}")
        formatted_results.append(f"Text:\n{result['text']}\n")  # Display full text
    return "\n".join(formatted_results)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_files():
    if 'files[]' not in request.files:
        return redirect(request.url)
    
    files = request.files.getlist('files[]')
    for file in files:
        if file:
            filename = file.filename
            file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(file_path)

    data = process_files(app.config['UPLOAD_FOLDER'])
    create_embeddings(data)

    return redirect(url_for('index'))

@app.route('/search', methods=['POST'])
def search_files():
    query = request.form['query']
    results = search(query)
    formatted_results = format_search_results(results)
    return render_template('index.html', query=query, results=formatted_results)

@app.route('/show_contents', methods=['GET'])
def show_contents():
    data = process_files(app.config['UPLOAD_FOLDER'])
    all_contents = "\n".join([item['text'] for item in data])
    return render_template('index.html', all_contents=all_contents)

if __name__ == '__main__':
    app.run(debug=True)