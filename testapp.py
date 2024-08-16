import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
from dotenv import load_dotenv
from txtai.embeddings import Embeddings
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from io import BytesIO
from docx import Document
from pptx import Presentation
import csv

load_dotenv()

txtai_embeddings = Embeddings({"path": "sentence-transformers/paraphrase-MiniLM-L3-v2", "content": True})

UPLOAD_DIR = "uploaded_files"
DOCS_DIR = "documentation"

if not os.path.exists(UPLOAD_DIR):
    os.makedirs(UPLOAD_DIR)

def save_uploaded_file(uploaded_file):
    with open(os.path.join(UPLOAD_DIR, uploaded_file.name), "wb") as f:
        f.write(uploaded_file.getvalue())

def get_existing_files(directory=UPLOAD_DIR):
    return [os.path.join(directory, f) for f in os.listdir(directory) if os.path.isfile(os.path.join(directory, f))]

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_docx_text(docx_docs):
    text = ""
    for docx in docx_docs:
        doc = Document(docx)
        for para in doc.paragraphs:
            text += para.text + "\n"
    return text

def get_pptx_text(pptx_docs):
    text = ""
    for pptx in pptx_docs:
        presentation = Presentation(pptx)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_txt_text(txt_docs):
    text = ""
    for txt_file in txt_docs:
        text += txt_file.read().decode('utf-8') + "\n"
    return text

def get_csv_text(csv_docs):
    text = ""
    for csv_file in csv_docs:
        try:
            decoded_content = csv_file.read().decode('utf-8', errors='ignore')
            csvreader = csv.reader(decoded_content.splitlines())
            for row in csvreader:
                text += ' '.join(row) + "\n"
        except Exception as e:
            st.error(f"Error reading CSV file: {e}")
    return text

def create_txtai_embeddings(text_chunks):
    data = [{"text": chunk} for chunk in text_chunks]
    txtai_embeddings.index(data)

def search_txtai(query):
    results = txtai_embeddings.search(query, limit=50)
    return results

def get_text_chunks(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = splitter.split_text(text)
    return chunks

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details.\n\n
    Context:\n{context}?\n
    Question:\n{question}\n
    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-pro")
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(llm=model, chain_type="stuff", prompt=prompt)
    return chain

def clear_chat_history():
    st.session_state.messages = [
        {"role": "assistant", "content": "You can ask questions about the provided documents."}]

def user_input(user_question):
    docs = search_txtai(user_question)

    chain = get_conversational_chain()

    context = "\n".join([doc["text"] for doc in docs])
    response = chain({"input_documents": docs, "context": context, "question": user_question}, return_only_outputs=True)

    return response['output_text']

def main():
    st.set_page_config(page_title="NYP Chatbot", page_icon="🖐", layout="wide")

    st.title("NYP Chatbot 🙋‍♂")

    st.title("Upload and Process Documents")

    uploaded_files = st.file_uploader("Upload your files (PDF, DOCX, PPTX, TXT, CSV)", accept_multiple_files=True, type=['pdf', 'docx', 'pptx', 'txt', 'csv'])

    if uploaded_files:
        for file in uploaded_files:
            save_uploaded_file(file)
    
    existing_files = get_existing_files(DOCS_DIR)
    all_text = ""
    pdf_files = [file for file in existing_files if file.endswith(".pdf")]
    docx_files = [file for file in existing_files if file.endswith(".docx")]
    pptx_files = [file for file in existing_files if file.endswith(".pptx")]
    txt_files = [open(file, "rb") for file in existing_files if file.endswith(".txt")]
    csv_files = [open(file, "rb") for file in existing_files if file.endswith(".csv")]

    if pdf_files:
        all_text += get_pdf_text(pdf_files)
    if docx_files:
        all_text += get_docx_text(docx_files)
    if pptx_files:
        all_text += get_pptx_text(pptx_files)
    if txt_files:
        all_text += get_txt_text(txt_files)
    if csv_files:
        all_text += get_csv_text(csv_files)

    if all_text:
        chunks = get_text_chunks(all_text)
        create_txtai_embeddings(chunks)

    st.write("Welcome to the chat!")
    st.sidebar.button('Clear Chat History', on_click=clear_chat_history)

    if "messages" not in st.session_state.keys():
        st.session_state.messages = [
            {"role": "assistant", "content": "You can ask questions about the provided documents."}]

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message['content'])

    if prompt := st.chat_input():
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)

        if st.session_state.messages[-1]["role"] != "assistant":
            with st.chat_message("assistant"):
                with st.spinner("Thinking..."):
                    response = user_input(prompt)
                    st.session_state.messages.append({"role": "assistant", "content": response})
                    st.markdown(response)

    with st.sidebar:
        st.title("Search and View Content")
        search_query = st.text_input("Enter search query")
        if st.button("Search"):
            if search_query:
                results = search_txtai(search_query)
                st.write("Search Results:")
                for result in results:
                    st.write(result["text"])
            else:
                st.write("Please enter a query to search.")

        if st.button("View Content"):
            if existing_files:
                st.write("Content of the uploaded files:")
                st.write(all_text)
            else:
                st.write("No files uploaded.")

if __name__ == "__main__":
    main()
